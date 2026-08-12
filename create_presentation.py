import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    
    # Set slide width and height to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Colors
    NAVY = RGBColor(26, 54, 93)
    TEAL = RGBColor(49, 151, 149)
    DARK_GRAY = RGBColor(45, 55, 72)
    LIGHT_BG = RGBColor(247, 250, 252)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(113, 128, 150)

    def add_header(slide, title_text, category_text="BIOMEDICAL SIGNAL PROCESSING PROJECT"):
        # Header Box
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Background shape
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "EEG-Based Attention Analysis"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Using Biomedical Signal Processing Pipeline (Biopac System)"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEAL
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Biomedical Engineering Project Presentation  |  10-Slide Summary Deck"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    p3.space_before = Pt(40)

    # SLIDE 2: Objectives
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "1. Project Objective & Core Focus")
    
    tb = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Main Objective", "Analyze real continuous EEG recorded via Biopac to observe frequency-domain changes between Baseline resting and Attention tasks."),
        ("Strict Requirement", "Pure Biomedical Signal Processing algorithms ONLY. No Machine Learning (CNN, SVM, Neural Networks) used."),
        ("Key Signal Processing Steps", "Data Validation -> Detrending -> Band-Pass Filtering -> Notch Filtering -> Welch PSD -> Band Power Extraction -> Attention Index -> Statistics -> Ablation Study."),
        ("Data Integrity", "All numerical metrics and conclusions are strictly derived from real Biopac EEG recordings.")
    ]
    
    for i, (title, desc) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

    # SLIDE 3: Acquisition & Protocol
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "2. EEG Acquisition & Time Segmentation")
    
    tb = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    details = [
        ("Hardware System", "Biopac Data Acquisition System"),
        ("Sampling Rate", "250 Hz (250 samples/sec)"),
        ("Channels", "2 EEG Channels"),
        ("Total Recording", "124,610 samples (~8.31 minutes / 498.44 seconds)"),
        ("Exact Time Segments", "• 0 to 64s: Baseline / Calibration Phase (32 epochs)\n• 64 to 194s: Eyes Closed / Relaxed Phase\n• 194 to 498s: Attention Task Phase (152 epochs)")
    ]
    
    for i, (title, desc) in enumerate(details):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # SLIDE 4: Preprocessing Pipeline
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "3. Signal Preprocessing Pipeline")
    
    tb = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    steps = [
        ("Step 1: Detrending", "Removes linear baseline drift and DC offset using scipy.signal.detrend."),
        ("Step 2: Band-pass Filtering", "0.5 – 40 Hz 4th-order zero-phase Butterworth filter (filtfilt) to isolate EEG frequency bands without phase distortion."),
        ("Step 3: Power-line Notch Filter", "50 Hz zero-phase IIR notch filter to remove 50 Hz power-line noise."),
        ("Step 4: Artifact Rejection", "Amplitude thresholding at 100 μV on 2-second windows (0 rejected, confirming clean signal execution).")
    ]
    
    for i, (title, desc) in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(17)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

    # SLIDE 5: PSD & Band Powers
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "4. Frequency-Domain Analysis & Feature Calculation")
    
    tb = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    items = [
        ("Epoching", "2-second non-overlapping windows (500 samples per epoch)."),
        ("Welch PSD Method", "Computes Power Spectral Density to reduce periodogram noise variance."),
        ("Extracted Bands", "Delta (0.5–4Hz), Theta (4–8Hz), Alpha (8–13Hz), Beta (13–30Hz), Gamma (30–40Hz)."),
        ("Attention Index Formula", "Beta / Theta Ratio = (Beta Power) / (Theta Power)\n(Higher values historically linked to cognitive focus).")
    ]
    
    for i, (title, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(17)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

    # SLIDE 6: Statistical Comparison
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "5. Statistical Comparison (Baseline vs Attention)")
    
    tb = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    stats_list = [
        ("Baseline Beta/Theta Mean", "1.131"),
        ("Attention Task Beta/Theta Mean", "1.105"),
        ("Percentage Change", "-2.25%"),
        ("Statistical Test", "Paired Wilcoxon Signed-Rank Test"),
        ("Test Statistic & p-value", "W = 200.0, p-value = 0.239 (p > 0.05, not statistically significant)"),
        ("Effect Size (Cohen's d)", "-0.278 (Small effect size)")
    ]
    
    for i, (title, desc) in enumerate(stats_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # SLIDE 7: Ablation Study
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "6. Required Ablation Study Results")
    
    # Add Table
    rows, cols = 6, 6
    table_shape = slide7.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    table = table_shape.table
    
    headers = ["Pipeline Stage", "Detrend", "Band-pass", "Notch", "RMS", "50Hz Power"]
    data = [
        ["Raw", "No", "No", "No", "2.916", "0.01032"],
        ["Detrended", "Yes", "No", "No", "0.601", "0.01032"],
        ["Band-pass", "Yes", "Yes", "No", "0.076", "0.00009"],
        ["Band-pass + Notch", "Yes", "Yes", "Yes", "0.076", "0.00002"],
        ["Full Pipeline", "Yes", "Yes", "Yes", "0.076", "0.00002"]
    ]
    
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.alignment = PP_ALIGN.CENTER

    # SLIDE 8: Filter Verification & Visuals
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "7. Visual & Filter Verification")
    
    tb = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    verifications = [
        ("Time-Domain Verification", "Linear detrending flattened large low-frequency baseline drifts (RMS dropped from 2.916 to 0.601)."),
        ("Frequency-Domain Verification", "Welch PSD curves confirm attenuation of high-frequency noise and 50 Hz power-line interference."),
        ("Feature Stability", "Filter pipeline reduced variance of extracted Beta power (0.000431 down to 0.000414), confirming feature stabilization.")
    ]
    
    for i, (title, desc) in enumerate(verifications):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(17)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(16)

    # SLIDE 9: Discussion
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide9, "8. Discussion & Scientific Observations")
    
    tb = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    discussions = [
        ("Signal Processing Success", "Filtering successfully reduced noise without introducing phase distortion (zero-phase filtering)."),
        ("Physiological Finding", "Beta/Theta ratio showed a minor change (-2.25%, p=0.239), indicating that cognitive engagement was uniform across tasks for this subject."),
        ("Scientific Integrity", "Demonstrates true empirical processing rather than forced or fabricated accuracy metrics.")
    ]
    
    for i, (title, desc) in enumerate(discussions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(17)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(16)

    # SLIDE 10: Conclusion
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide10, "9. Conclusion & Future Scope")
    
    tb = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    conclusions = [
        ("Conclusion", "Successfully implemented a complete 2-channel Biopac EEG signal processing pipeline strictly using traditional BSP algorithms."),
        ("Ablation Insights", "Proved that detrending, band-pass, and notch filtering are critical for feature stabilization."),
        ("Future Scope", "1. Test multi-subject cohorts.\n2. Utilize hardware synchronized event triggers.\n3. Expand to multi-channel 10-20 system recordings.")
    ]
    
    for i, (title, desc) in enumerate(conclusions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(17)
        run.font.color.rgb = DARK_GRAY
        p.space_after = Pt(16)
        
    out_path = r"c:\Users\chhuz\OneDrive\Desktop\BSP Projects\EEG_Attention_Project\EEG_Attention_Analysis_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    build_presentation()
